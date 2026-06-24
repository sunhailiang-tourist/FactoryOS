#!/usr/bin/env python3
"""Generate FactoryOS internal executive deck from v3.1 outline."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "准备" / "2026-06-16" / "FactoryOS-内部宣讲-v3.1.pptx"
ARCH = ROOT / "docs" / "文档" / "架构"

NAVY = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
GRAY = RGBColor(0x55, 0x55, 0x55)


def set_title(slide, text: str, subtitle: str | None = None):
    slide.shapes.title.text = text
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_bullets(slide, lines: list[str], left=0.6, top=1.6, width=12.0, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]], col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide, title)
    cols, rs = len(headers), len(rows) + 1
    tbl = slide.shapes.add_table(rs, cols, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.4 * rs)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
    return slide


def add_image_slide(prs, title: str, image_path: Path, caption: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide, title)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.4), width=Inches(12.3))
    if caption:
        add_bullets(slide, [caption], top=6.8, size=14)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "FactoryOS"
    s.placeholders[1].text = (
        "让工厂「会说话就能办事」，且「敢写库、能撤回」\n\n"
        "内部拍板 · 是什么 · 为什么 · 怎么干 · 多少钱 · 怎么长大\n"
        "2026-06 · v3.1"
    )

    # 2 Overlay
    add_table_slide(
        prs,
        "不换 ERP 的执行 OS——这就是 Overlay",
        ["✅ Overlay 是", "❌ 不是"],
        [
            ["接在 ERP/钉钉 上的执行入口 + 门禁", "替换金蝶/用友"],
            ["工人钉钉说/扫 → 主管确认 → 写 ERP", "把 ERP 抄进我们数仓"],
            ["Graph + Rule + Audit + Revert + 对账", "聊天机器人"],
            ["按 Pack 模块化卖、可复制", "每家从零定制项目"],
        ],
        [6.2, 6.2],
    )
    s = prs.slides[-1]
    add_bullets(
        s,
        [
            "工人/主管 → FactoryOS → 确认·审计·撤回",
            "              ↓ 写",
            "       客户 ERP（账本仍在这里）+ 钉钉（入口与待办）",
        ],
        top=5.2,
        size=14,
    )

    # 3 D1
    add_table_slide(
        prs,
        "先做生产的通病，不做全厂全链",
        ["原因", "说明"],
        [
            ["痛点 universal", "工人不爱用系统、主管怕写错、账实对不上"],
            ["价值可感知", "30 秒报工 vs 10 分钟填表"],
            ["风险可控", "只动产量写入，不动财务/autopay"],
            ["可验收", "D1 五项 + ≤90 天书面结案"],
        ],
        [3.5, 8.8],
    )
    add_bullets(
        prs.slides[-1],
        [
            "D1 五项：只读查询 → 受控写入 → 全链路审计 → 可 Revert → 每日对账",
            "哈森首期：ERP 报工 + 钉钉（无独立 MES）",
        ],
        top=5.0,
        size=16,
    )

    # 4 Overlay vs 数仓
    add_table_slide(
        prs,
        "先叠层，不搬库——两条路线的差别",
        ["", "Overlay（我们现在）", "数仓化/替换（不做主线）"],
        [
            ["卖什么", "敢写、可撤、有人用的执行", "数据汇聚、分析平台"],
            ["ERP 角色", "权威账本（L0）", "逐步变成被动库"],
            ["客户阻力", "低：「不换 ERP」", "高：IT/财务/金蝶渠道"],
            ["交付", "6～8 周内核 + 4～5 月首家", "常变 12 月+ 集成项目"],
            ["复制", "import Pack", "每家一套 ETL"],
        ],
        [2.2, 5.0, 5.2],
    )
    add_bullets(
        prs.slides[-1],
        [
            "档位1 Y1～ 控制面上浮（现在）| 档位2 Y2+ 读分析 | 档位3 选配数据池",
            "没有信任和结案，客户不会心甘情愿把 ERP 当数仓用",
        ],
        top=5.5,
        size=14,
    )

    # 5 下沉条件
    add_table_slide(
        prs,
        "数仓下沉的前提——信任先于搬迁",
        ["#", "条件"],
        [
            ["1", "≥1 家 D1 书面结案（Shadow、Revert、对账零 unexplained drift）"],
            ["2", "工人/主管日常只走 FactoryOS，Legacy 少人登录"],
            ["3", "客户书面采购 Layer 3 / 分析模块（单独合同）"],
            ["4", "对账证明：OS 记录 ↔ ERP 长期一致"],
        ],
        [0.8, 11.5],
    )
    add_bullets(
        prs.slides[-1],
        [
            "下沉 ≠ 把 ERP 全量抄到云上",
            "下沉 = 人不再点 ERP 界面 · 账本仍可留在 ERP · 我们持执行证明",
        ],
        top=5.2,
        size=15,
    )

    # 6 双线先锋
    add_table_slide(
        prs,
        "同一内核，两条先锋线——扩大可打市场",
        ["先锋线", "客户", "接法", "何时启动"],
        [
            ["A·有系统 Overlay", "哈森等（ERP+钉钉）", "conn-erp 写+读 + 钉钉", "内核 Gate 0 后立即"],
            ["B·无系统 Lite", "中小厂无 ERP/MES", "conn-mes-builtin 内置账本", "内核稳+哈森 Shadow 后（约 M4）"],
        ],
        [2.8, 2.5, 4.0, 3.0],
    )
    add_bullets(
        prs.slides[-1],
        [
            "同一套 OS（Graph/Rule/Revert/Audit）——只换 Connector 目标",
            "内核必须先硬（6～8 周）；B 线晚 1～2 月，不抢哈森 ERP 接口",
            "人力：2 人内核 → Gate 0 后 +1 接 A → M4 可选 +1 开 B",
        ],
        top=4.8,
        size=14,
    )

    # 7 路径总表
    add_table_slide(
        prs,
        "拍板后 12 个月怎么走",
        ["阶段", "时间", "做什么", "交付/Gate"],
        [
            ["0 拍板", "本周", "PPT 确认 + ADR 签字", "方向锁定"],
            ["1 内核", "W1～W8", "Graph/Rule/Exec/Audit/Revert；mock Connector", "Gate 0"],
            ["2 上云", "与 1 并行", "阿里云采购上线", "可部署"],
            ["3 哈森接入", "W7～W12", "ERP 报工写+读 + 钉钉 H5", "ERP+钉钉闭环"],
            ["4 哈森 D1", "+≤90 天", "Shadow → UAT → 结案", "首单 L0+L1"],
            ["5 B 线试点", "M4 起", "1 家无系统 B-Lite（可选）", "Starter-B"],
            ["6 复制", "M6～M12", "第 2～3 家 import Pack", "3～5 tenant"],
        ],
        [1.5, 1.5, 5.5, 3.8],
    )
    add_bullets(
        prs.slides[-1],
        ["关键数字：内核 6～8 周 · 首家 ~4～5 月 · 第二家 2～3 周 · 12 月初具规模"],
        top=6.5,
        size=14,
    )

    # 8 执行顺序
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(s, "怎么干：执行顺序")
    add_bullets(
        s,
        [
            "① 文档已齐 → ② 内核 Coding（AI Coding + AC 驱动）",
            "③ 云环境就绪 → ④ 哈森 ERP API + 钉钉应用",
            "⑤ Graph 工作坊 freeze → ⑥ Shadow ≥14 天 → ⑦ D1 签字",
            "⑧ export Package → ⑨ 第二家 import / B 线接入",
            "",
            "人力递进：",
            "· 内核：2 后端 | 收 ERP 文档",
            "· 哈森：+1 ERP/钉钉 +1 H5 | 1 人 Shadow/UAT",
            "· 复制/B 线：维持 2～3 研发 | +0.5～1 实施",
        ],
        top=1.5,
        size=17,
    )

    # 9 飞轮
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(s, "定制要变成 Pack，不是变成债务")
    add_bullets(
        s,
        [
            "客户现场定制（D2）→ ≥70% 沉淀为 Pack → export Package",
            "→ 下一家 import（2～3 周）→ Pack 目录越来越厚 = 飞轮",
            "",
            "场景定制：某厂特异逻辑 → Override/L2 Pack → ≥70% 入库",
            "数据飞轮：可卖的 Pack + 对账/Revert 模板 → Pack Registry",
            "铁律：无 export = 交付未完成；不改 execution 内核",
        ],
        top=1.5,
        size=17,
    )

    # 10 壁垒
    add_table_slide(
        prs,
        "壁垒不是模型，是机制 + 复利资产",
        ["阶", "时间", "沉淀什么", "壁垒"],
        [
            ["1", "Gate 0", "Graph/Rule/Revert/Audit 机制", "敢写库"],
            ["2", "首家 D1", "报工 Graph + ERP Connector + Package v1", "有人用、可结案"],
            ["3", "D2+复制", "扩展 Pack + import 4～8 周→2～3 周", "复制速度"],
            ["4", "12 月", "Pack 目录 + 多 tenant License", "规模毛利"],
            ["5", "选配", "Layer 3 分析/数仓加深", "客户信任后加价"],
        ],
        [0.6, 1.2, 5.5, 4.5],
    )
    add_bullets(
        prs.slides[-1],
        [
            "四层壁垒：产品机制 | 工程纪律（11 条红线）| 交付纪律 | 复利资产",
            "外挂 Agent 1～2 周能 demo；这套 OS 6 个月起，难在机制+纪律+Pack",
        ],
        top=5.5,
        size=14,
    )

    # 11 增长
    add_table_slide(
        prs,
        "两条先锋线汇合后，客户数怎么涨",
        ["里程碑", "商业信号"],
        [
            ["首家 D1 结案", "证明能收钱、能交付"],
            ["第二家 ≤3 周", "证明不是项目制"],
            ["B 线 1 家试点", "证明 TAM 扩大（无系统厂也能进）"],
            ["12 月 3～5 tenant", "初具规模（原 Y3 目标前移）"],
        ],
        [4.5, 7.8],
    )
    add_bullets(
        prs.slides[-1],
        [
            "A 有 ERP：哈森 → 同行业 import | B 无系统：B-Lite 试点",
            "定价：L0 订阅 + L1 Pack + 实施",
        ],
        top=5.0,
        size=15,
    )

    # 12 成本
    add_table_slide(
        prs,
        "多少钱、多少人（一页说清）",
        ["类别", "年费/节奏", "备注"],
        [
            ["软硬件合计", "~5.4～9.1 万/年", "已含阿里云（约 3.7～5.5 万）"],
            ["其中·云", "（已含在上行）", "~3,700～4,600/月；一次齐套，只升配"],
            ["其中·Cursor+API+SSL", "（已含在上行）", "Cursor $60/人、研发 API、证书"],
            ["月均（2～3 人 Cursor）", "~5,100～7,400/月", "软硬件总账，非「云+另加」"],
            ["人力·内核期", "2 研发", "6～8 周"],
            ["人力·哈森", "+2～3", "Gate 0 后"],
            ["人力·B 线", "+1（可选 M4）", "不抢 A 线"],
        ],
        [3.5, 3.0, 5.8],
    )

    # 13 拍板
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(s, "请老板拍板")
    add_bullets(
        s,
        [
            "□ 1. 方向   制造业 AI 执行 OS（Overlay），不换 ERP",
            "□ 2. 切入   哈森：生产报工 · ERP 写+读 · 钉钉入口",
            "□ 3. 双线   有系统 Overlay 先行；无系统 B-Lite 内核稳后并行",
            "□ 4. 节奏   内核 6～8 周 → 首家 ~4～5 月 → 12 月初具规模",
            "□ 5. 资源   软硬件 ~5.4～9.1 万/年（含阿里云）+ 内核 2 人起",
            "□ 6. 配合   灯塔厂 + ERP 报工 API + 钉钉 + Graph 冻结",
            "□ 7. 纪律   D1 五项 · Shadow · 无 export 不算交付完成",
            "",
            "其他期望：_________________",
            "签字：_________________",
        ],
        top=1.4,
        size=18,
    )

    # Appendix A one-pager
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(s, "附录 A｜一页纸（可打印）")
    add_bullets(
        s,
        [
            "是什么：Overlay = 钉钉入口 + ERP 账本不动 + 敢写可撤",
            "先做什么：生产报工通病（D1 五项）· 哈森 ERP+钉钉",
            "为什么不数仓：先信任与结案；数仓是档位3选配",
            "双线：有系统 Overlay + 无系统 B-Lite（同内核）",
            "路径：拍板→内核6～8周→哈森4～5月→复制→12月3～5家",
            "资产：Pack + Package export · 定制≥70%入库",
            "成本：软硬件 ~5.4～9.1万/年（含云）· 人力 2人起",
        ],
        top=1.5,
        size=17,
    )

    # Appendix E
    add_table_slide(
        prs,
        "附录 E-1｜两本账（备查）",
        ["账本", "含什么", "年费粗算"],
        [
            ["软硬件合计", "云 + Cursor + 研发 API + 域名 SSL", "~5.4～9.1 万/年"],
            ["人力", "研发 + 实施工资", "另计（2 人内核起）"],
        ],
        [2.5, 5.5, 4.3],
    )
    add_bullets(prs.slides[-1], ["云 ≈ 房子的租金 | Cursor/API ≈ 盖房子的工具"], top=4.5, size=16)

    add_table_slide(
        prs,
        "附录 E-2｜测试域 vs 生产域（备查）",
        ["", "测试域", "生产域"],
        [
            ["干什么", "代码仓库·CI·staging·API 文档", "只跑 FactoryOS 业务"],
            ["核心组件", "GitLab+Runner+staging+Redoc", "SLB+API+RDS+Redis+OSS"],
            ["数据库", "RDS 测试（可重建）", "RDS 生产（15 天备份）"],
            ["谁用", "研发", "工人/主管（钉钉）"],
            ["禁止", "—", "GitLab、YApi、自建 PG"],
        ],
        [2.0, 5.0, 5.3],
    )
    add_bullets(
        prs.slides[-1],
        [
            "测试四类：① GitLab ② Runner+staging ③ Redoc ④ RDS 测试",
        ],
        top=5.8,
        size=14,
    )

    add_table_slide(
        prs,
        "附录 E-3｜非云研发工具（备查）",
        ["工具", "跑在哪", "干什么", "账本"],
        [
            ["Cursor Pro+", "研发本机", "AI 写内核代码", "$60/人/月"],
            ["LiteLLM", "API 出站", "意图/填槽/DSL", "研发期→12；生产→tenant"],
            ["ASR/OCR/VLM", "API 出站", "说/拍报工联调", "CI mock $0"],
            ["钉钉开放平台", "生产出站", "H5/待办/通知", "低/免费"],
        ],
        [2.2, 2.0, 4.5, 3.5],
    )

    # Appendix F with images
    add_image_slide(
        prs,
        "附录 F-1｜系统架构图",
        ARCH / "系统架构图.png",
        "L3 界面 → L2 Agent → L0 内核 → L1 Connector → Data-L0 客户账本 | 哈森：ERP+钉钉",
    )
    add_image_slide(
        prs,
        "附录 F-2｜技术架构图",
        ARCH / "技术架构图.png",
        "仅 agent_orchestrator 用 LLM；唯一写路径：execution → connector → Legacy",
    )
    add_image_slide(
        prs,
        "附录 F-3｜数据架构图",
        ARCH / "数据架构图.png",
        "L0 客户 ERP | L1 运行数据→RDS | L2 缓存 | L3 选配 | 媒体→OSS",
    )

    # Appendix G
    add_table_slide(
        prs,
        "附录 G｜研发全流程 × 资源（备查）",
        ["步骤", "时间", "人", "云/工具", "交付"],
        [
            ["拍板", "本周", "负责人", "—", "方向锁定"],
            ["内核", "W1～W8", "2 后端", "Cursor·GitLab·staging·mock AI", "Gate 0"],
            ["上云", "并行", "0.5", "按 10 采购", "测试+生产就绪"],
            ["哈森接入", "Gate 0 后", "+2", "生产→ERP/钉钉", "报工闭环"],
            ["Shadow→D1", "≤90d", "1 实施", "生产 RDS/OSS", "首单结案"],
            ["复制/B-Lite", "M4+", "+0～1", "同云多 tenant", "import 2～3 周"],
        ],
        [1.8, 1.5, 1.5, 4.0, 3.5],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Generated {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
